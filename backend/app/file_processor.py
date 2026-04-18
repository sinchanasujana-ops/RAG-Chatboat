import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
import tempfile
import os
from fastapi import UploadFile

# Point pytesseract to tesseract executable (Windows)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Supported extensions
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx", ".jpg", ".jpeg", ".png"}


def get_extension(filename: str) -> str:
    return os.path.splitext(filename)[-1].lower()


def is_supported(filename: str) -> bool:
    return get_extension(filename) in SUPPORTED_EXTENSIONS


# ── Extractors ───────────────────────────────────────

def extract_from_pdf(file_path: str) -> list[dict]:
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append({
                    "page_number": i + 1,
                    "text": text.strip()
                })
    return pages


def extract_from_docx(file_path: str) -> list[dict]:
    doc = Document(file_path)
    full_text = "\n".join(
        para.text for para in doc.paragraphs if para.text.strip()
    )
    if not full_text.strip():
        return []
    return [{"page_number": 1, "text": full_text.strip()}]


def extract_from_txt(file_path: str) -> list[dict]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()
    if not text:
        return []
    return [{"page_number": 1, "text": text}]


def extract_from_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            pages.append({
                "page_number": i + 1,
                "text": "\n".join(texts)
            })
    return pages


def extract_from_image(file_path: str) -> list[dict]:
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image).strip()
    if not text:
        return []
    return [{"page_number": 1, "text": text}]


# ── Main Dispatcher ──────────────────────────────────

def extract_text(file_path: str, filename: str) -> list[dict]:
    """
    Detects file type by extension and extracts text accordingly.
    Returns list of {page_number, text} dicts.
    """
    ext = get_extension(filename)

    if ext == ".pdf":
        return extract_from_pdf(file_path)
    elif ext == ".docx":
        return extract_from_docx(file_path)
    elif ext == ".txt":
        return extract_from_txt(file_path)
    elif ext == ".pptx":
        return extract_from_pptx(file_path)
    elif ext in {".jpg", ".jpeg", ".png"}:
        return extract_from_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ── File Save Helper ─────────────────────────────────

async def save_upload_file_temporarily(upload_file: UploadFile) -> str:
    suffix = get_extension(upload_file.filename)
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        content = await upload_file.read()
        tmp.write(content)
        return tmp.name